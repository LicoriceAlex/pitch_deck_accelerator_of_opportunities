import os

from PIL import Image

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor

from django.conf import settings

from .investments import make_req


class PitchDeckGenerator:
    def __init__(self, content,
                 template_name, logo_path,
                 presentation_path) -> None:
        self.content = content
        self.logo_path = logo_path
        self.presentation_path = presentation_path
        self.template = os.path.join(settings.BASE_DIR,
                                     'presentations',
                                     'pres_templates',
                                     f'{template_name}.pptx')

    def create_pptx(self):
        # from pptx import Presentation
        prs = Presentation(self.template)

        # Слайд 1
        slide1 = prs.slides[0]
        slide1.shapes.title.text = self.content.get('name')
        image_path = self.logo_path
        left = Inches(13.33/2 - 3/2)  # X-координата
        top = Inches(7.5/2 + 0.25)   # Y-координата
        width = Inches(3)  # Ширина изображения
        height = Inches(3)  # Высота изображения
        pic = slide1.shapes.add_picture(image_path, left, top, width, height)

        # Слайд 2
        slide2 = prs.slides[1]
        slide2.placeholders[1].text = self.content.get('problem')
        left = Inches(13.33/2 + 3/2 + 1)  # X-координата
        top = Inches(7.5/2 - 1.25)   # Y-координата
        width = Inches(3)  # Ширина изображения
        height = Inches(3)  # Высота изображения
        pic2 = slide2.shapes.add_picture(image_path, left, top, width, height)

        # Слайд 3
        slide3 = prs.slides[2]
        slide3.placeholders[1].text = self.content.get('description')
        left = Inches(13.33/2 - 3/2 - 4)  # X-координата
        top = Inches(7.5/2 - 1.25)   # Y-координата
        width = Inches(3)  # Ширина изображения
        height = Inches(3)  # Высота изображения
        pic3 = slide3.shapes.add_picture(image_path, left, top, width, height)

        # Слайд 4
        slide4 = prs.slides[3]
        slide4.placeholders[1].text = self.content.get('solution')
        left = Inches(13.33/2 + 3/2 + 1)  # X-координата
        top = Inches(7.5/2 - 1.25)   # Y-координата
        width = Inches(3)  # Ширина изображения
        height = Inches(3)  # Высота изображения
        pic4 = slide4.shapes.add_picture(image_path, left, top, width, height)

        # Слайд 5
        slide5 = prs.slides[4]
        slide5.placeholders[1].text = 'TAM\n\nSAM\n\nSOM'

        # Слайд 6
        # Таблица бизнес модель
        # slide6 = prs.slides[5]
        # x, y, cx, cy = Inches(3), Inches(2), Inches(8), Inches(6)
        # shape = slide6.shapes.add_table(9, 2, x, y, cx, cy)

        # Слайд 7
        slide7 = prs.slides[6]
        slide7.placeholders[1].text = 'Здесь будет текст'
        left = Inches(13.33/2 + 3/2 + 1)  # X-координата
        top = Inches(7.5/2 - 1.25)   # Y-координата
        width = Inches(3)  # Ширина изображения
        height = Inches(3)  # Высота изображения
        pic7 = slide7.shapes.add_picture(image_path, left, top, width, height)


        # Слайд 8
        slide8 = prs.slides[7]
        slide8.placeholders[1].text = self.content.get('team')

        # Слайд 9
        slide9 = prs.slides[8]
        left = Inches(2)  # X-координата
        top = Inches(2)   # Y-координата
        width = Inches(13.33 - 2 - 2)  # Ширина изображения
        height = Inches(7.5 - 2 - 2)  # Высота изображения
        desired_investments = self.content.get('desired_investments')
        investment_round = self.content.get('investment_round')
        investments = self.content.get('investments')
        fig_path = make_req(desired_investments, investment_round, investments)
        # fig_path = some_func() -> path
        fig = slide9.shapes.add_picture(fig_path, left, top, width, height)

        # Слайд 10
        slide10 = prs.slides[9]
        slide10.placeholders[1].text = self.content.get('roadmap')

        # Слайд 11
        slide11 = prs.slides[10]
        slide11.placeholders[1].text = self.content.get('contact_information')

        prs.save(self.presentation_path)

    # def delete_first_slides(self, presentation, first_num):
    #     slide_ids = reversed(range(first_num))
    #     for slide_id in slide_ids:
    #         if slide_id < len(presentation.slides):
    #             xml_slides = presentation.slides._sldIdLst
    #             slides = list(xml_slides)
    #             xml_slides.remove(slides[slide_id])

    # def create_ppt(self):
    #     prs = Presentation(self.template)

    #     content_slide_layout = prs.slide_layouts[1]

    #     for slide_content in self.slides_content:
    #         slide = prs.slides.add_slide(content_slide_layout)

    #         for placeholder in slide.placeholders:
    #             if placeholder.placeholder_format.type == 1:  # Title
    #                 placeholder.left = Inches(1)
    #                 placeholder.top = Inches(0.5)
    #                 placeholder.width = Inches(5)
    #                 placeholder.height = Inches(1)
    #                 placeholder.text = slide_content['title']
    #                 for paragraph in placeholder.text_frame.paragraphs:
    #                     for run in paragraph.runs:
    #                         run.font.name = 'Arial'
    #                         run.font.size = Pt(25)
    #             elif placeholder.placeholder_format.type == 4:  # Content

    #                 placeholder.left = Inches(1)
    #                 placeholder.top = Inches(1.25)
    #                 placeholder.width = Inches(7)
    #                 placeholder.height = Inches(4)
    #                 placeholder.text = slide_content['content']
    #                 for paragraph in placeholder.text_frame.paragraphs:
    #                     for run in paragraph.runs:
    #                         run.font.name = 'Arial'
    #                         run.font.size = Pt(11)

    #         if self.logo_path:
    #             # add the image at the specified position
    #             slide_width = Cm(16)
    #             slide_height = Cm(9)

    #             im = Image.open(self.logo_path)
    #             width, height = im.size

    #             image_width, image_height = Cm(width / 700), Cm(height / 700)

    #             left = slide_width - image_width + Cm(8)   # calculate left position
    #             top = slide_height - image_height - Cm(5)  # calculate top position

    #             slide.shapes.add_picture(self.logo_path,
    #                                      left, top,
    #                                      width=image_width,
    #                                      height=image_height)

    #     # Delete the first twelve slides after all new slides have been added
    #     self.delete_first_slides(prs, 12)

    #     # Save the presentation
    #     prs.save(self.presentation_path)